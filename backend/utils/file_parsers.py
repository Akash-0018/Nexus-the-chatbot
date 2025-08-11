import os
import tempfile
from typing import Dict, Any
import logging

# PDF parsing
try:
    import PyPDF2
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Excel parsing
try:
    import pandas as pd
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

# Word document parsing
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PowerPoint parsing
try:
    from pptx import Presentation
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

# Image processing
try:
    from PIL import Image
    import pytesseract
    import cv2
    import numpy as np
    IMAGE_AVAILABLE = True
except ImportError:
    IMAGE_AVAILABLE = False

class FileParser:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def parse_file(self, filepath: str, file_extension: str, original_filename: str) -> Dict[str, Any]:
        """Parse file content based on extension"""
        
        try:
            if file_extension == 'pdf':
                return self.parse_pdf(filepath)
            elif file_extension in ['xlsx', 'xls']:
                return self.parse_excel(filepath)
            elif file_extension in ['docx', 'doc']:
                return self.parse_docx(filepath)
            elif file_extension in ['pptx', 'ppt']:
                return self.parse_ppt(filepath)
            elif file_extension in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
                return self.parse_image(filepath)
            elif file_extension == 'txt':
                return self.parse_text(filepath)
            else:
                return {'error': f'Unsupported file type: {file_extension}'}
        
        except Exception as e:
            self.logger.error(f"Error parsing {file_extension} file: {str(e)}")
            return {'error': f'Failed to parse {file_extension} file: {str(e)}'}
    
    def parse_pdf(self, filepath: str) -> Dict[str, Any]:
        """Extract text from PDF files"""
        if not PDF_AVAILABLE:
            return {'error': 'PDF parsing libraries not available'}
        
        try:
            reader = PdfReader(filepath)
            text_content = ""
            metadata = {
                'total_pages': len(reader.pages),
                'title': reader.metadata.title if reader.metadata else None,
                'author': reader.metadata.author if reader.metadata else None
            }
            
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                text_content += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
            
            return {
                'type': 'pdf',
                'text_content': text_content.strip(),
                'metadata': metadata,
                'summary': f"PDF document with {len(reader.pages)} pages extracted successfully"
            }
            
        except Exception as e:
            return {'error': f'PDF parsing failed: {str(e)}'}
    
    def parse_excel(self, filepath: str) -> Dict[str, Any]:
        """Extract data from Excel files"""
        if not EXCEL_AVAILABLE:
            return {'error': 'Excel parsing libraries not available'}
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(filepath)
            sheets_data = {}
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(filepath, sheet_name=sheet_name)
                sheets_data[sheet_name] = {
                    'data': df.to_dict('records'),
                    'shape': df.shape,
                    'columns': df.columns.tolist(),
                    'summary': f"Sheet '{sheet_name}' has {df.shape[0]} rows and {df.shape[1]} columns"
                }
            
            return {
                'type': 'excel',
                'sheets': sheets_data,
                'total_sheets': len(excel_file.sheet_names),
                'sheet_names': excel_file.sheet_names,
                'summary': f"Excel file with {len(excel_file.sheet_names)} sheets processed successfully"
            }
            
        except Exception as e:
            return {'error': f'Excel parsing failed: {str(e)}'}
    
    def parse_docx(self, filepath: str) -> Dict[str, Any]:
        """Extract text from Word documents"""
        if not DOCX_AVAILABLE:
            return {'error': 'Word document parsing libraries not available'}
        
        try:
            doc = Document(filepath)
            
            # Extract paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            
            # Extract tables
            tables_data = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables_data.append(table_data)
            
            full_text = '\n'.join(paragraphs)
            
            return {
                'type': 'docx',
                'text_content': full_text,
                'paragraphs': paragraphs,
                'tables': tables_data,
                'stats': {
                    'paragraph_count': len(paragraphs),
                    'table_count': len(tables_data),
                    'word_count': len(full_text.split())
                },
                'summary': f"Word document with {len(paragraphs)} paragraphs and {len(tables_data)} tables extracted"
            }
            
        except Exception as e:
            return {'error': f'Word document parsing failed: {str(e)}'}
    
    def parse_ppt(self, filepath: str) -> Dict[str, Any]:
        """Extract text from PowerPoint presentations"""
        if not PPT_AVAILABLE:
            return {'error': 'PowerPoint parsing libraries not available'}
        
        try:
            prs = Presentation(filepath)
            slides_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                slides_content.append({
                    'slide_number': slide_num + 1,
                    'content': slide_text,
                    'text': '\n'.join(slide_text)
                })
            
            all_text = '\n\n'.join([slide['text'] for slide in slides_content])
            
            return {
                'type': 'pptx',
                'text_content': all_text,
                'slides': slides_content,
                'total_slides': len(slides_content),
                'summary': f"PowerPoint presentation with {len(slides_content)} slides extracted"
            }
            
        except Exception as e:
            return {'error': f'PowerPoint parsing failed: {str(e)}'}
    
    def parse_image(self, filepath: str) -> Dict[str, Any]:
        """Extract text from images using OCR"""
        if not IMAGE_AVAILABLE:
            return {'error': 'Image processing libraries not available'}
        
        try:
            # Open and analyze image
            image = Image.open(filepath)
            
            # Get image metadata
            metadata = {
                'format': image.format,
                'mode': image.mode,
                'size': image.size,
                'width': image.width,
                'height': image.height
            }
            
            # Try OCR text extraction
            extracted_text = ""
            try:
                extracted_text = pytesseract.image_to_string(image)
            except Exception as ocr_error:
                extracted_text = f"OCR failed: {str(ocr_error)}"
            
            return {
                'type': 'image',
                'metadata': metadata,
                'extracted_text': extracted_text.strip(),
                'has_text': bool(extracted_text.strip()),
                'summary': f"Image ({image.width}x{image.height}) processed successfully"
            }
            
        except Exception as e:
            return {'error': f'Image processing failed: {str(e)}'}
    
    def parse_text(self, filepath: str) -> Dict[str, Any]:
        """Read plain text files"""
        try:
            with open(filepath, 'r', encoding='utf-8') as file:
                content = file.read()
            
            return {
                'type': 'text',
                'text_content': content,
                'stats': {
                    'character_count': len(content),
                    'word_count': len(content.split()),
                    'line_count': len(content.splitlines())
                },
                'summary': f"Text file with {len(content.split())} words processed"
            }
            
        except Exception as e:
            return {'error': f'Text file parsing failed: {str(e)}'}
