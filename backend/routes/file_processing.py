from flask import Blueprint, request, jsonify, current_app
from flask_cors import cross_origin
from werkzeug.utils import secure_filename
import os
import uuid
from datetime import datetime
import tempfile
import logging

file_bp = Blueprint('files', __name__)
logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = {
    'pdf', 'txt', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt',
    'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp'
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def parse_file_content(filepath, file_extension, filename):
    """Parse uploaded file content - Your enhanced version"""
    try:
        if file_extension == 'txt':
            # Handle multiple encodings for text files
            for encoding in ['utf-8', 'utf-16', 'utf-16-le', 'latin-1']:
                try:
                    with open(filepath, 'r', encoding=encoding, errors='ignore') as f:
                        content = f.read().replace('\x00', '').strip()
                    if content:
                        return {
                            'type': 'text',
                            'text_content': content,
                            'word_count': len(content.split()),
                            'summary': f'Text file with {len(content.split())} words processed successfully'
                        }
                except:
                    continue
        
        elif file_extension == 'pdf':
            try:
                from pypdf import PdfReader
                reader = PdfReader(filepath)
                text_pages = []
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ''
                    if page_text.strip():
                        text_pages.append(f"--- Page {i+1} ---\n{page_text}")
                
                full_text = '\n\n'.join(text_pages)
                return {
                    'type': 'pdf',
                    'text_content': full_text,
                    'page_count': len(reader.pages),
                    'word_count': len(full_text.split()),
                    'summary': f'PDF with {len(reader.pages)} pages processed successfully'
                }
            except ImportError:
                return {
                    'type': 'pdf',
                    'summary': f'PDF "{filename}" uploaded (install pypdf for text extraction)',
                    'requires': 'pypdf'
                }
        
        elif file_extension == 'docx':
            try:
                from docx import Document
                doc = Document(filepath)
                paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                full_text = '\n\n'.join(paragraphs)
                
                return {
                    'type': 'docx',
                    'text_content': full_text,
                    'paragraph_count': len(paragraphs),
                    'word_count': len(full_text.split()),
                    'summary': f'Word document with {len(paragraphs)} paragraphs processed successfully'
                }
            except ImportError:
                return {
                    'type': 'docx',
                    'summary': f'Word document "{filename}" uploaded (install python-docx)',
                    'requires': 'python-docx'
                }
        
        elif file_extension in ['xlsx', 'xls']:
            try:
                import pandas as pd
                df = pd.read_excel(filepath)
                return {
                    'type': 'excel',
                    'data_preview': df.head(5).to_dict('records'),
                    'shape': df.shape,
                    'columns': df.columns.tolist(),
                    'summary': f'Excel file with {df.shape[0]} rows and {df.shape[1]} columns processed'
                }
            except ImportError:
                return {
                    'type': 'excel',
                    'summary': f'Excel "{filename}" uploaded (install pandas)',
                    'requires': 'pandas'
                }
        
        # Add PowerPoint support (missing from minimal version)
        elif file_extension in ['pptx', 'ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                slides_content = []
                
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        slides_content.append({
                            'slide_number': slide_num + 1,
                            'content': slide_text,
                            'text': '\n'.join(slide_text)
                        })
                
                all_text = '\n\n'.join([slide['text'] for slide in slides_content])
                
                return {
                    'type': 'pptx',
                    'text_content': all_text,
                    'slides': slides_content,
                    'total_slides': len(slides_content),
                    'summary': f'PowerPoint with {len(slides_content)} slides processed successfully'
                }
            except ImportError:
                return {
                    'type': 'pptx',
                    'summary': f'PowerPoint "{filename}" uploaded (install python-pptx)',
                    'requires': 'python-pptx'
                }
        
        # Add Image OCR support (missing from minimal version)
        elif file_extension in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
            try:
                from PIL import Image
                image = Image.open(filepath)
                
                extracted_text = ""
                try:
                    import pytesseract
                    extracted_text = pytesseract.image_to_string(image).strip()
                except ImportError:
                    extracted_text = "OCR not available (install pytesseract)"
                except Exception as e:
                    extracted_text = f"OCR failed: {str(e)}"
                
                return {
                    'type': 'image',
                    'metadata': {
                        'format': image.format,
                        'size': image.size,
                        'mode': image.mode
                    },
                    'extracted_text': extracted_text,
                    'has_text': bool(extracted_text and "OCR" not in extracted_text),
                    'summary': f'Image "{filename}" ({image.size[0]}x{image.size[1]}) processed successfully'
                }
            except ImportError:
                return {
                    'type': 'image',
                    'summary': f'Image "{filename}" received (processing requires Pillow)',
                    'requires': 'Pillow'
                }
        
        else:
            return {
                'type': file_extension,
                'summary': f'File "{filename}" uploaded successfully'
            }
    
    except Exception as e:
        logger.error(f"File parsing error: {e}")
        return {
            'error': f'Error processing {filename}: {str(e)}',
            'summary': f'File processing failed'
        }

# **ONLY CHANGE: Add credentials support to CORS decorators**
@file_bp.route('/upload', methods=['POST', 'OPTIONS'])
@cross_origin(supports_credentials=True)  # ✅ ONLY CHANGE: Add credentials support
def upload_file():
    """Handle file upload with proper CORS support"""
    
    # Handle preflight OPTIONS request
    if request.method == 'OPTIONS':
        response = jsonify({'status': 'ok'})
        response.headers['Access-Control-Allow-Credentials'] = 'true'  # ✅ ONLY CHANGE: Add this
        return response, 200
    
    try:
        logger.info("File upload request received")
        
        # Check if file is present
        if 'file' not in request.files:
            return jsonify({
                'success': False, 
                'error': 'No file provided'
            }), 400
        
        file = request.files['file']
        
        # Validate file
        if not file.filename:
            return jsonify({
                'success': False, 
                'error': 'No file selected'
            }), 400
        
        if not allowed_file(file.filename):
            return jsonify({
                'success': False,
                'error': f'File type not supported. Allowed: {", ".join(ALLOWED_EXTENSIONS)}'
            }), 400
        
        # Check file size
        if request.content_length and request.content_length > 16 * 1024 * 1024:
            return jsonify({
                'success': False,
                'error': 'File too large. Maximum size: 16MB'
            }), 413
        
        # Generate unique identifiers
        file_id = str(uuid.uuid4())
        filename = secure_filename(file.filename)
        file_extension = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
        
        # Process file in temporary location
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
            file.save(temp_file.name)
            temp_path = temp_file.name
        
        logger.info(f"Processing file: {filename}")
        
        # Parse file content
        parsed_content = parse_file_content(temp_path, file_extension, filename)
        
        # Cleanup temporary file
        try:
            os.unlink(temp_path)
        except:
            pass
        
        logger.info(f"File processed successfully: {filename}")
        
        return jsonify({
            'success': True,
            'file_id': file_id,
            'filename': filename,
            'file_type': file_extension,
            'content': parsed_content,
            'timestamp': datetime.now().isoformat(),
            'message': f'File "{filename}" processed successfully by Nexus!'
        })
        
    except Exception as e:
        logger.error(f"File upload error: {e}")
        return jsonify({
            'success': False,
            'error': f'Upload failed: {str(e)}'
        }), 500

@file_bp.route('/test', methods=['GET'])
@cross_origin(supports_credentials=True)  # ✅ ONLY CHANGE: Add credentials support
def test_files():
    """Test file processing system"""
    return jsonify({
        'status': 'File processing system operational',
        'cors_enabled': True,
        'credentials_supported': True,  # ✅ ONLY CHANGE: Add this indicator
        'supported_formats': list(ALLOWED_EXTENSIONS),
        'max_file_size': '16MB',
        'success': True,
        'timestamp': datetime.now().isoformat()
    })
